#!/usr/bin/env python3
"""
Excel to PowerPoint Generator
Generates McKinsey-grade presentations from application data.

Dependencies (conservative, bank-approved):
    - openpyxl: Excel file reading
    - python-pptx: PowerPoint generation

Usage:
    python excel_to_ppt.py <input_excel.xlsx> [output.pptx]
"""

import sys
import os
from collections import defaultdict
from datetime import datetime

# Conservative imports - widely used in enterprise environments
try:
    from openpyxl import load_workbook
except ImportError:
    print("ERROR: openpyxl not installed. Run: pip install openpyxl")
    sys.exit(1)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx")
    sys.exit(1)


# ==============================================================================
# CONFIGURATION - McKinsey Color Palette
# ==============================================================================
class Colors:
    """McKinsey-inspired professional color palette"""
    DARK_BLUE = RGBColor(0, 51, 102)      # Primary headers
    MEDIUM_BLUE = RGBColor(0, 102, 153)   # Secondary headers
    LIGHT_BLUE = RGBColor(204, 229, 255)  # Table header background
    DARK_GRAY = RGBColor(51, 51, 51)      # Body text
    LIGHT_GRAY = RGBColor(242, 242, 242)  # Alternating rows
    WHITE = RGBColor(255, 255, 255)
    BLACK = RGBColor(0, 0, 0)
    ACCENT_GREEN = RGBColor(0, 128, 0)    # For highlights


class Config:
    """Presentation configuration"""
    SLIDE_WIDTH = Inches(13.333)  # Widescreen 16:9
    SLIDE_HEIGHT = Inches(7.5)
    MARGIN_LEFT = Inches(0.5)
    MARGIN_TOP = Inches(1.2)
    MARGIN_RIGHT = Inches(0.5)
    MARGIN_BOTTOM = Inches(0.5)
    
    # Font sizes
    TITLE_FONT_SIZE = Pt(28)
    SUBTITLE_FONT_SIZE = Pt(18)
    HEADER_FONT_SIZE = Pt(11)
    BODY_FONT_SIZE = Pt(9)
    FOOTER_FONT_SIZE = Pt(8)
    
    # Table settings
    MAX_ROWS_PER_SLIDE = 12
    TABLE_ROW_HEIGHT = Inches(0.35)


# ==============================================================================
# EXCEL DATA READER
# ==============================================================================
class ExcelReader:
    """Reads and structures data from the Excel file"""
    
    REQUIRED_COLUMNS = [
        'Application Name',
        'Application Hosting Country',
        'LOBT',
        'DATA DOMAIN FULL',
        'SUB-DATA DOMAIN',
        'SUB-SUB-DATA DOMAIN',
        'ACTIVITY',
        'Host Application Code'
    ]
    
    def __init__(self, filepath):
        self.filepath = filepath
        self.data = []
        self.column_map = {}
        
    def read(self):
        """Read Excel file and return structured data"""
        print(f"Reading Excel file: {self.filepath}")
        
        wb = load_workbook(filename=self.filepath, read_only=True, data_only=True)
        ws = wb.active
        
        # Get headers from first row
        headers = []
        for cell in ws[1]:
            headers.append(str(cell.value).strip() if cell.value else '')
        
        # Map column indices
        self._map_columns(headers)
        
        # Read data rows
        for row_idx, row in enumerate(ws.iter_rows(min_row=2), start=2):
            row_data = self._extract_row(row)
            if row_data:
                self.data.append(row_data)
        
        wb.close()
        print(f"Read {len(self.data)} data rows")
        return self.data
    
    def _map_columns(self, headers):
        """Map required columns to their indices"""
        for col_name in self.REQUIRED_COLUMNS:
            found = False
            for idx, header in enumerate(headers):
                if self._normalize(header) == self._normalize(col_name):
                    self.column_map[col_name] = idx
                    found = True
                    break
            if not found:
                # Try partial match
                for idx, header in enumerate(headers):
                    if self._normalize(col_name) in self._normalize(header):
                        self.column_map[col_name] = idx
                        found = True
                        break
        
        missing = [c for c in self.REQUIRED_COLUMNS if c not in self.column_map]
        if missing:
            print(f"WARNING: Could not find columns: {missing}")
            print(f"Available columns: {headers}")
    
    def _normalize(self, text):
        """Normalize column name for matching"""
        return text.lower().replace(' ', '').replace('-', '').replace('_', '')
    
    def _extract_row(self, row):
        """Extract data from a row"""
        def get_value(col_name):
            if col_name in self.column_map:
                idx = self.column_map[col_name]
                if idx < len(row):
                    val = row[idx].value
                    return str(val).strip() if val else ''
            return ''
        
        data_domain = get_value('DATA DOMAIN FULL')
        if not data_domain:
            return None
            
        return {
            'data_domain': data_domain,
            'sub_domain': get_value('SUB-DATA DOMAIN'),
            'sub_sub_domain': get_value('SUB-SUB-DATA DOMAIN'),
            'activity': get_value('ACTIVITY'),
            'app_code': get_value('Host Application Code'),
            'app_name': get_value('Application Name'),
            'country': get_value('Application Hosting Country'),
            'lobt': get_value('LOBT')
        }
    
    def get_hierarchical_data(self):
        """Organize data hierarchically by domain structure"""
        hierarchy = defaultdict(lambda: defaultdict(lambda: defaultdict(lambda: defaultdict(list))))
        
        for row in self.data:
            domain = row['data_domain'] or 'Unspecified'
            sub = row['sub_domain'] or 'Unspecified'
            sub_sub = row['sub_sub_domain'] or 'Unspecified'
            activity = row['activity'] or 'Unspecified'
            
            app_info = {
                'app_code': row['app_code'],
                'app_name': row['app_name'],
                'country': row['country'],
                'lobt': row['lobt']
            }
            
            hierarchy[domain][sub][sub_sub][activity].append(app_info)
        
        return hierarchy


# ==============================================================================
# POWERPOINT GENERATOR
# ==============================================================================
class PPTGenerator:
    """Generates McKinsey-grade PowerPoint presentations"""
    
    def __init__(self, output_path):
        self.output_path = output_path
        self.prs = Presentation()
        self.prs.slide_width = Config.SLIDE_WIDTH
        self.prs.slide_height = Config.SLIDE_HEIGHT
        self.slide_count = 0
        
    def create_presentation(self, hierarchy, source_file):
        """Create the full presentation"""
        print("Generating PowerPoint presentation...")
        
        # Title slide
        self._add_title_slide(source_file)
        
        # Executive summary
        self._add_summary_slide(hierarchy)
        
        # Domain slides
        for domain in sorted(hierarchy.keys()):
            self._add_domain_section(domain, hierarchy[domain])
        
        # Save
        self.prs.save(self.output_path)
        print(f"Presentation saved: {self.output_path}")
        print(f"Total slides: {self.slide_count}")
        
    def _add_title_slide(self, source_file):
        """Add professional title slide"""
        slide = self._add_blank_slide()
        
        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(12.3), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Application Data Domain Mapping"
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = Colors.DARK_BLUE
        p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.6), Inches(12.3), Inches(0.5)
        )
        tf = subtitle_box.text_frame
        p = tf.paragraphs[0]
        p.text = "Application Portfolio by Data Domain and Activity"
        p.font.size = Pt(20)
        p.font.color.rgb = Colors.MEDIUM_BLUE
        p.alignment = PP_ALIGN.CENTER
        
        # Date and source
        date_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(5), Inches(12.3), Inches(0.5)
        )
        tf = date_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}"
        p.font.size = Pt(12)
        p.font.color.rgb = Colors.DARK_GRAY
        p.alignment = PP_ALIGN.CENTER
        
        # Add decorative line
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(4), Inches(4.2), Inches(5.3), Pt(3)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = Colors.DARK_BLUE
        line.line.fill.background()
        
    def _add_summary_slide(self, hierarchy):
        """Add executive summary slide"""
        slide = self._add_blank_slide()
        self._add_slide_title(slide, "Executive Summary")
        
        # Calculate statistics
        total_domains = len(hierarchy)
        total_sub_domains = sum(len(subs) for subs in hierarchy.values())
        total_apps = set()
        for domain in hierarchy.values():
            for sub in domain.values():
                for sub_sub in sub.values():
                    for activity in sub_sub.values():
                        for app in activity:
                            if app['app_code']:
                                total_apps.add(app['app_code'])
        
        # Summary text
        summary_box = slide.shapes.add_textbox(
            Config.MARGIN_LEFT, Inches(1.5), Inches(12), Inches(1)
        )
        tf = summary_box.text_frame
        tf.word_wrap = True
        
        stats = [
            f"• Total Data Domains: {total_domains}",
            f"• Total Sub-Data Domains: {total_sub_domains}",
            f"• Total Unique Applications: {len(total_apps)}",
        ]
        
        for i, stat in enumerate(stats):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = stat
            p.font.size = Pt(16)
            p.font.color.rgb = Colors.DARK_GRAY
            p.space_after = Pt(12)
        
        # Domain overview table
        self._add_domain_overview_table(slide, hierarchy)
        
    def _add_domain_overview_table(self, slide, hierarchy):
        """Add overview table showing domain counts"""
        rows = len(hierarchy) + 1
        cols = 3
        
        table = slide.shapes.add_table(
            rows, cols,
            Config.MARGIN_LEFT, Inches(3.5),
            Inches(8), Inches(0.35 * min(rows, 8))
        ).table
        
        # Headers
        headers = ['Data Domain', 'Sub-Domains', 'Applications']
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            self._format_header_cell(cell)
        
        # Data rows
        for row_idx, domain in enumerate(sorted(hierarchy.keys()), start=1):
            if row_idx >= rows:
                break
            
            sub_domains = hierarchy[domain]
            app_count = set()
            for sub in sub_domains.values():
                for sub_sub in sub.values():
                    for activity in sub_sub.values():
                        for app in activity:
                            if app['app_code']:
                                app_count.add(app['app_code'])
            
            data = [
                domain[:50] + '...' if len(domain) > 50 else domain,
                str(len(sub_domains)),
                str(len(app_count))
            ]
            
            for col_idx, value in enumerate(data):
                cell = table.cell(row_idx, col_idx)
                cell.text = value
                self._format_body_cell(cell, row_idx % 2 == 0)
        
        # Set column widths
        table.columns[0].width = Inches(5)
        table.columns[1].width = Inches(1.5)
        table.columns[2].width = Inches(1.5)
        
    def _add_domain_section(self, domain, sub_domains):
        """Add slides for a data domain"""
        # Domain divider slide
        self._add_divider_slide(domain)
        
        # For each sub-domain
        for sub_domain in sorted(sub_domains.keys()):
            sub_sub_domains = sub_domains[sub_domain]
            
            # Collect all applications for this sub-domain
            all_apps = []
            for sub_sub in sorted(sub_sub_domains.keys()):
                activities = sub_sub_domains[sub_sub]
                for activity in sorted(activities.keys()):
                    apps = activities[activity]
                    for app in apps:
                        all_apps.append({
                            'sub_sub_domain': sub_sub,
                            'activity': activity,
                            **app
                        })
            
            # Create slides with tables
            self._add_sub_domain_slides(domain, sub_domain, all_apps)
    
    def _add_divider_slide(self, domain):
        """Add a section divider slide"""
        slide = self._add_blank_slide()
        
        # Background shape
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0), Inches(2.8), Config.SLIDE_WIDTH, Inches(2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = Colors.DARK_BLUE
        shape.line.fill.background()
        
        # Domain title
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(3.2), Inches(12.3), Inches(1)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = domain
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = Colors.WHITE
        p.alignment = PP_ALIGN.CENTER
        
    def _add_sub_domain_slides(self, domain, sub_domain, apps):
        """Add content slides for a sub-domain"""
        # Paginate if needed
        page_size = Config.MAX_ROWS_PER_SLIDE
        pages = [apps[i:i + page_size] for i in range(0, len(apps), page_size)]
        
        if not pages:
            pages = [[]]
        
        for page_num, page_apps in enumerate(pages, start=1):
            slide = self._add_blank_slide()
            
            # Title with domain context
            title = f"{domain}"
            subtitle = f"Sub-Domain: {sub_domain}"
            if len(pages) > 1:
                subtitle += f" (Page {page_num}/{len(pages)})"
            
            self._add_slide_title(slide, title, subtitle)
            
            # Add application table
            self._add_app_table(slide, page_apps)
    
    def _add_app_table(self, slide, apps):
        """Add application data table"""
        if not apps:
            # No data message
            msg_box = slide.shapes.add_textbox(
                Config.MARGIN_LEFT, Inches(3), Inches(12), Inches(1)
            )
            tf = msg_box.text_frame
            p = tf.paragraphs[0]
            p.text = "No applications mapped to this sub-domain"
            p.font.size = Pt(14)
            p.font.color.rgb = Colors.DARK_GRAY
            p.font.italic = True
            return
        
        rows = len(apps) + 1
        cols = 6
        
        table = slide.shapes.add_table(
            rows, cols,
            Config.MARGIN_LEFT, Config.MARGIN_TOP + Inches(0.5),
            Inches(12.3), Config.TABLE_ROW_HEIGHT * rows
        ).table
        
        # Headers
        headers = ['Sub-Sub Domain', 'Activity', 'App Code', 'Application Name', 'Country', 'LOBT']
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = header
            self._format_header_cell(cell)
        
        # Data rows
        for row_idx, app in enumerate(apps, start=1):
            data = [
                self._truncate(app.get('sub_sub_domain', ''), 25),
                self._truncate(app.get('activity', ''), 25),
                self._truncate(app.get('app_code', ''), 15),
                self._truncate(app.get('app_name', ''), 35),
                self._truncate(app.get('country', ''), 10),
                self._truncate(app.get('lobt', ''), 15)
            ]
            
            for col_idx, value in enumerate(data):
                cell = table.cell(row_idx, col_idx)
                cell.text = value
                self._format_body_cell(cell, row_idx % 2 == 0)
        
        # Set column widths
        table.columns[0].width = Inches(2.2)
        table.columns[1].width = Inches(2.2)
        table.columns[2].width = Inches(1.3)
        table.columns[3].width = Inches(3.5)
        table.columns[4].width = Inches(1.1)
        table.columns[5].width = Inches(2.0)
    
    def _add_blank_slide(self):
        """Add a blank slide"""
        self.slide_count += 1
        blank_layout = self.prs.slide_layouts[6]  # Blank layout
        slide = self.prs.slides.add_slide(blank_layout)
        self._add_footer(slide)
        return slide
    
    def _add_slide_title(self, slide, title, subtitle=None):
        """Add title to slide"""
        # Title
        title_box = slide.shapes.add_textbox(
            Config.MARGIN_LEFT, Inches(0.3), Inches(12), Inches(0.6)
        )
        tf = title_box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Config.TITLE_FONT_SIZE
        p.font.bold = True
        p.font.color.rgb = Colors.DARK_BLUE
        
        if subtitle:
            subtitle_box = slide.shapes.add_textbox(
                Config.MARGIN_LEFT, Inches(0.85), Inches(12), Inches(0.4)
            )
            tf = subtitle_box.text_frame
            p = tf.paragraphs[0]
            p.text = subtitle
            p.font.size = Config.SUBTITLE_FONT_SIZE
            p.font.color.rgb = Colors.MEDIUM_BLUE
        
        # Title underline
        line = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Config.MARGIN_LEFT, Inches(1.15), Inches(12.3), Pt(2)
        )
        line.fill.solid()
        line.fill.fore_color.rgb = Colors.DARK_BLUE
        line.line.fill.background()
    
    def _add_footer(self, slide):
        """Add footer to slide"""
        footer_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.3)
        )
        tf = footer_box.text_frame
        p = tf.paragraphs[0]
        p.text = f"Application Data Domain Mapping | Confidential | Page {self.slide_count}"
        p.font.size = Config.FOOTER_FONT_SIZE
        p.font.color.rgb = Colors.DARK_GRAY
        p.alignment = PP_ALIGN.RIGHT
    
    def _format_header_cell(self, cell):
        """Format table header cell"""
        cell.fill.solid()
        cell.fill.fore_color.rgb = Colors.DARK_BLUE
        
        para = cell.text_frame.paragraphs[0]
        para.font.size = Config.HEADER_FONT_SIZE
        para.font.bold = True
        para.font.color.rgb = Colors.WHITE
        para.alignment = PP_ALIGN.CENTER
        
        cell.text_frame.paragraphs[0].font.size = Config.HEADER_FONT_SIZE
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    def _format_body_cell(self, cell, alternate=False):
        """Format table body cell"""
        if alternate:
            cell.fill.solid()
            cell.fill.fore_color.rgb = Colors.LIGHT_GRAY
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = Colors.WHITE
        
        para = cell.text_frame.paragraphs[0]
        para.font.size = Config.BODY_FONT_SIZE
        para.font.color.rgb = Colors.DARK_GRAY
        para.alignment = PP_ALIGN.LEFT
        
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    def _truncate(self, text, max_len):
        """Truncate text with ellipsis"""
        if not text:
            return ''
        text = str(text)
        if len(text) > max_len:
            return text[:max_len-3] + '...'
        return text


# ==============================================================================
# MAIN
# ==============================================================================
def main():
    """Main entry point"""
    print("=" * 60)
    print("Excel to PowerPoint Generator")
    print("=" * 60)
    
    # Parse arguments
    if len(sys.argv) < 2:
        print("\nUsage: python excel_to_ppt.py <input.xlsx> [output.pptx]")
        print("\nExample:")
        print("  python excel_to_ppt.py applications.xlsx output.pptx")
        sys.exit(1)
    
    input_file = sys.argv[1]
    
    if not os.path.exists(input_file):
        print(f"ERROR: Input file not found: {input_file}")
        sys.exit(1)
    
    # Output file
    if len(sys.argv) >= 3:
        output_file = sys.argv[2]
    else:
        base = os.path.splitext(input_file)[0]
        output_file = f"{base}_presentation.pptx"
    
    print(f"\nInput:  {input_file}")
    print(f"Output: {output_file}\n")
    
    # Read Excel
    reader = ExcelReader(input_file)
    reader.read()
    hierarchy = reader.get_hierarchical_data()
    
    if not hierarchy:
        print("ERROR: No data found in Excel file")
        sys.exit(1)
    
    # Generate PowerPoint
    generator = PPTGenerator(output_file)
    generator.create_presentation(hierarchy, input_file)
    
    print("\n" + "=" * 60)
    print("Generation complete!")
    print("=" * 60)


if __name__ == '__main__':
    main()
